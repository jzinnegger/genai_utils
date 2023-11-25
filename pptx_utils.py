
# Visualize the elements in a PowerPoint template for inspection
# Applicable to any PowerPoint template with multiple master layout
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
import pandas as pd
import json
from pptx.oxml.xmlchemy import OxmlElement

def _draw_box_around_placeholder(slide, placeholder):
    """
    Draws a box with no fill and a border around a given placeholder on a slide.
    """
    left, top, width, height = placeholder.left, placeholder.top, placeholder.width, placeholder.height
    box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    box.fill.background()  # Setting no fill for the box
    box.line.color.rgb = RGBColor(0, 0, 0)  # Black line color
    box.line.width = Pt(1)  # Line thickness


def layout_inspection(template_path, output_path=None):
    """
    Scans a given PowerPoint template for all masters and their layouts and placeholders,
    and optionally saves the layout inspection slide.

    Args:
    template_path (str): Path to the PowerPoint template file.
    output_path (str, optional): Path to save the output PowerPoint file. If None, the file is not saved.
    """
    prs = Presentation(template_path)
    data = []

    # Iterate through all master slides and their layouts
    for master_index, slide_master in enumerate(prs.slide_masters):
        for layout_index, slide_layout in enumerate(slide_master.slide_layouts):
            # Add a slide with the current layout
            slide = prs.slides.add_slide(slide_layout)

            # Set the title (if present) to include the master and layout index
            if slide.shapes.title:
                slide.shapes.title.text = f"Master: {master_index}, Layout: {layout_index} (Placeholder Index: 0)"

            # Iterate through all placeholders in the slide
            for shape in slide.placeholders:
                phf = shape.placeholder_format
                placeholder_index = phf.idx
                # Gather data for DataFrame
                data.append([master_index, layout_index, shape.name, f"{phf.type}", placeholder_index])
                # Set text and draw a box around each placeholder
                if placeholder_index != 0:  # Skip the title placeholder
                    shape.text = f"Placeholder Index: {placeholder_index}; Shape Name: {shape.name}; Placeholder Type: {phf.type}"
                    _draw_box_around_placeholder(slide, shape)

    # Create DataFrame from the collected data
    df = pd.DataFrame(data, columns=['Master Index', 'Layout Index', 'Shape Name', 'Placeholder Type', 'Placeholder Index'])

    # Save the presentation if output_path is provided
    if output_path:
        prs.save(output_path)
    
    return df



def change_table_style(shape, style_name='Light Style 1', accent='Accent 1'):
    """
        Styles ids are from Microsoft documentation 
        "https://learn.microsoft.com/en-us/previous-versions/office/developer/office-2010/hh273476(v=office.14)?redirectedfrom=MSDN"
        see also discussion on github "https://github.com/scanny/python-pptx/issues/27"
        Changes the style of a table in a given shape.
        style_name (str): see values below
        accent (str): see values below
            No Style: No Grid, Table Grid
            Themed Style 1: Accent 1, Accent 2, Accent 3, Accent 4, Accent 5, Accent 6
            Themed Style 2: Accent 1, Accent 2, Accent 3, Accent 4, Accent 5, Accent 6
            Light Style 1: Accent 0, Accent 1, Accent 2, Accent 3, Accent 4, Accent 5, Accent 6
            Light Style 2: Accent 0, Accent 1, Accent 2, Accent 3, Accent 4, Accent 5, Accent 6
            Light Style 3: Accent 0, Accent 1, Accent 2, Accent 3, Accent 4, Accent 5, Accent 6
            Medium Style 1: Accent 0, Accent 1, Accent 2, Accent 3, Accent 4, Accent 5, Accent 6
            Medium Style 2: Accent 0, Accent 1, Accent 2, Accent 3, Accent 4, Accent 5, Accent 6
            Medium Style 3: Accent 0, Accent 1, Accent 2, Accent 3, Accent 4, Accent 5, Accent 6
            Medium Style 4: Accent 0, Accent 1, Accent 2, Accent 3, Accent 4, Accent 5, Accent 6
            Dark Style 1: Accent 0, Accent 1, Accent 2, Accent 3, Accent 4, Accent 5, Accent 6
            Dark Style 2: Accent 0, Accent 1/Accent 2, Accent 3/Accent 4, Accent 5/Accent 6
    """
    # Load style data from JSON file
    with open('style_data.json', 'r') as file:
        table_style_dict = json.load(file)

    style_id = table_style_dict.get(style_name, {}).get(accent)


    if style_id and shape is not None and hasattr(shape, '_element') and shape._element.graphic.graphicData.tbl is not None:
        tbl = shape._element.graphic.graphicData.tbl

        # Find or create the tblPr element
        tblPr = tbl.xpath('./a:tblPr')
        if tblPr:
            tblPr = tblPr[0]
        else:
            tblPr = OxmlElement('a:tblPr')
            tbl.insert(0, tblPr)

        # Find or create the tableStyleId element
        tableStyleId_elem = tblPr.xpath('./a:tableStyleId')
        if tableStyleId_elem:
            tableStyleId_elem = tableStyleId_elem[0]
        else:
            tableStyleId_elem = OxmlElement('a:tableStyleId')
            tblPr.append(tableStyleId_elem)

        tableStyleId_elem.text = style_id  # Set the style ID


def insert_slide_with_table(
    table_df: pd.DataFrame, 
    title_text: str, 
    output_path: str, 
    template_path: str, 
    layout: int = 2, 
    table_plc_idx: int = 23, 
    message_text: str = None, 
    message_plc_idx: int = 22
) -> None:
    """
    Inserts a slide with a table into a PowerPoint presentation.

    :param table_df: A pandas DataFrame containing the content for the table.
    :param title_text: Text for the title of the slide.
    :param output_path: Path to save the output PowerPoint file.
    :param template_path: Path to the PowerPoint template file.
    :param layout: Index of the slide layout to use.
    :param table_plc_idx: Placeholder index for the table.
    :param message_text: Text for the message on the slide, if any.
    :param message_plc_idx: Placeholder index for the message.
    :return: None
    """

    prs = Presentation(template_path)

    # Add a slide with layout
    slide_layout = prs.slide_layouts[layout]
    slide = prs.slides.add_slide(slide_layout)

    # Set the title
    title_placeholder = slide.placeholders[0]
    title_placeholder.text = title_text

    # Check if the message placeholder exists and message text is provided
    if message_plc_idx < len(slide.placeholders):
        message_placeholder = slide.placeholders[message_plc_idx]
        if message_placeholder.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and message_text:
            message_placeholder.text = message_text

    # Add a table to the table placeholder
    table_placeholder = slide.placeholders[table_plc_idx]
    rows, cols = table_df.shape[0] + 1, table_df.shape[1] + 1
    shape = table_placeholder.insert_table(rows, cols)
    table = shape.table

    # Set column names and row headers
    for col in range(1, cols):
        table.cell(0, col).text = str(table_df.columns[col - 1])

    for row in range(1, rows):
        table.cell(row, 0).text = str(table_df.index[row - 1])
        for col in range(1, cols):
            table.cell(row, col).text = str(table_df.iloc[row - 1, col - 1])

    # Set the table style
    change_table_style(shape)

    # Save the presentation
    prs.save(output_path)